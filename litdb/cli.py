"""CLI for litdb

The main command is litdb. There are subcommands for the actions.
"""

import click
from rich import print as richprint
import os
from sentence_transformers import SentenceTransformer
import bs4
import requests
import datetime
import pathlib

from tqdm import tqdm
import datetime
import dateparser
import json
import numpy as np
import ollama
import time
import webbrowser
from docx import Document
from pptx import Presentation
import nbformat
from nbconvert import MarkdownExporter
    
from . import root, CONFIG, config
from .db import get_db, add_source, add_work, add_author, update_filter, add_bibtex
from .openalex import get_data, get_text
from .pdf import add_pdf

db = get_db()

@click.group()
def cli():
    """Group command for litdb."""
    pass
             

#################
# Add functions #
#################

@cli.command()
@click.argument('sources', nargs=-1)
@click.option('--references', is_flag=True, help='Add references too.')
@click.option('--related', is_flag=True, help='Add related too.')
@click.option('--citing', is_flag=True, help='Add citing too.')
def add(sources, references=False, citing=False, related=False):
    """Add WIDS to the db.

    SOURCES can be one or more of a doi or orcid, a pdf path, a url, bibtex
    file, or other kind of file assumed to be text.

    These are one time additions. 

    """
    for source in tqdm(sources):

        # a work
        if source.startswith('10.') or 'doi.org' in source:
            if source.startswith('10.'):
                source = f'https://doi.org/{source}'
            add_work(source, references, citing, related)

        # works from an author
        elif 'orcid' in source:
            add_author(source)

        # a bibtex file
        elif source.endswith('.bib'):
            add_bibtex(source)
            
        # pdf
        elif source.endswith('.pdf'):
            add_pdf(source)

        # docx
        elif source.endswith('.docx'):            
            doc = Document(source)
            add_source(source, '\n'.join([para.text for para in doc.paragraphs]))

        # pptx
        elif source.endswith('.pptx'):            
            prs = Presentation(source)
            text = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text.append(shape.text)
            add_source(source, '\n'.join(text))

        # html
        elif source.endswith('.html'):
            with open(source) as f:
                text = f.read()
            soup = bs4.BeautifulSoup(text, features="lxml")
            add_source(source, soup.get_text())

        # a url
        elif source.startswith('http'):
            soup = bs4.BeautifulSoup(requests.get(source).text)
            add_source(source, soup.get_text())

        # ipynb
        elif source.endswith('.ipynb'):
            with open(source, 'r', encoding='utf-8') as f:
                notebook = nbformat.read(f, as_version=4)

            # Create a Markdown exporter
            markdown_exporter = MarkdownExporter()

            # Convert the notebook to Markdown
            (body, resources) = markdown_exporter.from_notebook_node(notebook)

            add_source(source, body)

        # assume it is text
        else:            
            with open(source) as f:
                text = f.read()
            add_source(source, text)
            


@cli.command()
@click.argument('sources', nargs=-1)
def index(sources):
    """Index the directories in SOURCES.
    SOURCES is a list of directories.
    """
    
    for directory in sources:
        directory = pathlib.Path(directory).resolve()
        for fname in directory.rglob('*'):            
            # for f in files:            
            if fname.suffix in ['.pdf', '.docx', '.pptx', '.org', '.md', '.html', '.bib', '.ipynb']:
                fname = str(fname)
                    
                # skip files we already have
                if db.execute('''select source from sources where source = ?''', (fname,)).fetchone():
                    continue
                    
                try:
                    print(f'Adding {fname}')
                    add([fname])
                    # add seems to call a SystemExit, so this line doesn't get run. I don't know why this happens here.
                    richprint(f'Added {fname}')
                # I don't know why this gets called, but it does, and I catch it so we keep going
                except SystemExit:
                    pass
                except:
                    richprint('Still something worng')
                    import sys
                    richprint(sys.exc_info())

        last_updated = datetime.date.today().strftime('%Y-%m-%d')

        directory = str(directory) # we need strings for the db
        if db.execute('''select path from directories where path = ?''', (directory,)).fetchone():
            print(f'Updating {directory}')
            db.execute('''update directories set last_updated = ? where path = ?''',
                       (last_updated, directory))
        else:
            print(f'Inserting {directory}: {last_updated}')
            db.execute('''insert into directories(path, last_updated) values (?, ?)''', (directory, last_updated))

        db.commit()


@cli.command()
def reindex():
    """Reindex saved directories."""
    for directory, in db.execute('''select path from directories''').fetchall():
        print(f'Reindexing {directory}')
        index([directory])



##########
# Review #
##########
@cli.command()
@click.option('-s', '--since',  default='1 week ago')
@click.option('-f', '--format', '_format',  default='org')
def review(since, _format):
    """Review new entries added SINCE. This should be something dateparser can handle.

    The default format is org. Other formats might be supported in the future.
    """
    c = db.execute('''select source, text, extra from sources where date(date_added) > ?''',
                   (dateparser.parse(since).strftime("%Y-%m-%d"),)).fetchall()
    for source, text, extra in c:
        if _format == 'org':
            data = json.loads(extra) or {}
            print(f'''* {source}
:PROPERTIES:
:CITED_BY_COUNT: {data.get('cited_by_count', 0)}
:END:

{text}
''')
                   

#############
# Searching #
#############
      
@cli.command()
@click.argument('query', nargs=-1)
@click.option('-n',  default=3)
def vsearch(query, n=3):
    query = ' '.join(query)
    model = SentenceTransformer(config['embedding']['model'])
    emb = model.encode([query]).astype(np.float32).tobytes()
    c = db.execute('''select sources.source, sources.text, vector_distance_cos(?, embedding) from vector_top_k('embedding_idx', ?, ?)
    join sources on sources.rowid = id''',
    (emb, emb, n))
    for i, row in enumerate(c.fetchall()):
        source, text, similarity = row
        try:
            richprint(f'{i + 1:2d}. ({similarity:1.2f}) {text}\n\n')
        except:
            # rich.errors.MarkupError
            print(f'{i + 1:2d}. ({similarity:1.2f}) {text}\n\n')

        
@cli.command()
@click.argument('query', nargs=-1)
@click.option('-n',  default=3)
def fulltext(query, n):
    """Perform a fulltext search on litdb.
    """
    query = ' '.join(query)
    
    for source, text in db.execute('''select source, snippet(fulltext, 1, '', '', '', 16)
    from fulltext
    where text match ? order by rank limit ?''',
    (query, n)).fetchall():
        richprint(f"[link]{source}[/link]")
        richprint(text + '\n')

        
# Adapted from https://www.arsturn.com/blog/understanding-ollamas-embedding-models
@cli.command()
@click.argument('prompt', nargs=-1)
def gpt(prompt):
    """Run an ollama query with PROMPT.
    """
    t0 = time.time()
    prompt = ' '.join(prompt)
    model = SentenceTransformer(config['embedding']['model'])
    emb = model.encode([prompt]).astype(np.float32).tobytes()
    richprint(f'It took {time.time() - t0:1.1f} sec to embed the prompt')
    t0 = time.time()
    data = db.execute('''select sources.text from vector_top_k('embedding_idx', ?, 3) join sources on sources.rowid = id''',
    (emb,)).fetchall()
    richprint(f'It took  {time.time() - t0:1.1f} sec to get the top three docs')
    t0 = time.time()
    output = ollama.generate(model="llama2", prompt=f"Using data: {data}. Respond to the prompt: {prompt}")
    richprint(output['response'])
    richprint(f'It took  {time.time() - t0:1.1f} sec to generate and richprint the response.')

    richprint('The text was generated using these references')
    for i, result in enumerate(data):
        richprint(f'{i:2d}. {result}\n')        


@cli.command()
@click.argument('source')
@click.option('-n',  default=3)
def similar(source, n=3):
    emb, = db.execute('''select embedding from sources where source = ?''', (source,)).fetchone()

    # print starting at index 1, the first item is always the source.
    for i, row in enumerate(db.execute('''select sources.source, sources.text from vector_top_k('embedding_idx', ?, ?) join sources on sources.rowid = id''',
                                       # we do n + 1 because the first entry is always the source
                                       (emb, n + 1)).fetchall()[1:]):
        source, text = row
        richprint(f'{i:2d}. {source}\n{text}\n')
        
    
###########
# Filters #
###########

@cli.command()
@click.argument('filter')
@click.option('-d', '--description')
def add_filter(filter, description=None):
    """Add an OpenAlex FILTER. 
    """
    db.execute('insert into queries(filter, description) values (?, ?)',
               (filter, description))
    db.commit()

    
@cli.command()
@click.argument('filter')
def rm_filter(filter):
    """Remove an OpenAlex FILTER. 
    """
    db.execute('delete from queries where filter = ?',
               (filter,))
    db.commit()


@cli.command()
def update_filters():
    """Update litdb using a filter with works from a created date.
    """
    filters = db.execute('''select filter, last_updated from queries''')
    for f, last_updated in filters.fetchall():
        update_filter(f, last_updated)


@cli.command()        
def list_filters():
    """List the filters.
    """
    filters = db.execute('''select rowid, filter, description, last_updated from queries''')
    for rowid, f, description, last_updated in filters.fetchall():
        richprint(f'{rowid:3d}. {description or "None":30s} : {f} ({last_updated})')


######################
# OpenAlex searching #
######################

@cli.command()
@click.argument('query')
@click.option('-e', '--endpoint', default='works')
def openalex(query, endpoint='works'):
    """Run an openalex query on FILTER.

    ENDPOINT should be one of works, authors, or another entity.
    
    This does not add anything to your database. It is to help you find starting points.

    To search text:
    litdb openalex "default.search:circular polymer"

    To find a journal id
    litdb openalex -e sources "display_name.search:Digital Discovery"
    """
    url = f'https://api.openalex.org/{endpoint}'

    params={'email': config['openalex']['email'],
            'api_key': config['openalex'].get('api_key'),
            'filter': query}
    resp = requests.get(url, params)
    
    data = resp.json()
  
    print(f'Found {data["meta"]["count"]} results.')
    for result in data['results']:
        try:
            richprint(get_text(result))
            print()
        except:
            richprint(f'{result["id"]}: {result["display_name"]}')
    

########################################
# Convenience functions to add filters #
########################################
        
@cli.command()
@click.argument('name', nargs=-1)
def author_search(name):
    """Search OpenAlex for name.
    Uses the autocomplete endpoint to find an author's orcid.
    """
    auname = ' '.join(name)
    
    url = 'https://api.openalex.org/autocomplete/authors'

    from .openalex import get_data

    data = get_data(url,
                    params={'q': auname})

    for result in data['results']:
        richprint(f'- {result["display_name"]}\n  {result["hint"]} {result["external_id"]}\n\n')


@cli.command()
@click.argument('orcid')
@click.option('-r', '--remove', is_flag=True, help='remove')
def follow(orcid, remove=False):
    """Add a filter to follow orcid.
    """

    if not orcid.startswith('http'):
        orcid = f'https://orcid.org/{orcid}'

    # Seems like we should get their articles first.
    add_author(orcid)

    f = f'author.orcid:{orcid}'

    if remove:
        c = db.execute('''delete from queries where  filter = ?''',
                   (f,))
        db.commit()
        richprint(f'{c.rowcount} rows removed')
        return
        
    url = f'https://api.openalex.org/authors/{orcid}'
    data = get_data(url)
    name = data['display_name']  
    
    db.execute('''insert or ignore into queries(filter, description) values (?, ?)''',
               (f, name))
    richprint(f'Following {name}: {orcid}')
    db.commit()


@cli.command()
@click.argument('query', nargs=-1)
@click.option('-r', '--remove', is_flag=True, help='remove')
def watch(query, remove=False):
    """Setup a watch on query.
    QUERY: string, a filter for openalex.
    """
    
    # First, we should make sure the filter is valid
    query = ' '.join(query)

    if remove:
        c = db.execute('''delete from queries where filter = ?''', (query,))
        db.commit()
        richprint(f'{c.rowcount} rows removed')
        return
    
    url = 'https://api.openalex.org/works'

    data = get_data(url, params={'filter': query})
    if len(data['results']) == 0:
        richprint(f"Sorry, {query} does not seem valid.")

    if remove:
        c = db.execute('''delete from queries where filter = ?''', (query,))        
        richprint(f'Deleted {c.rowcount} rows')
        db.commit()
    else:
        c = db.execute('''insert or ignore into queries(filter, description) values (?, ?)''',
                       (query,))        
        richprint(f'Added {c.rowcount} rows')
        db.commit()
        richprint(f'Watching {query}')
    

@cli.command()
@click.argument('doi')
@click.option('-r', '--remove', is_flag=True, help='remove')
def citing(doi, remove=False):
    """Setup a watch for articles that cite doi.
    """
           
    url = 'https://api.openalex.org/works'

    # We need an OpenAlex id
    f = f'doi:{doi}'

    data = get_data(url, params={'filter': f})
    if len(data['results']) == 0:
        richprint(f"Sorry, {doi} does not seem valid.")

    wid = data['results'][0]['id']

    if remove:
        c = db.execute('''delete from queries where filter = ?''',
                   (f'cites:{wid}',))
        db.commit()
        richprint(f'Deleted {c.rowcount} rows')
    else:
        c = db.execute('''insert or ignore into queries(filter, description) values (?, ?)''',
               (f'cites:{wid}', f'Citing papers for {doi}'))
        
        db.commit()
        richprint(f'Added {c.rowcount} rows')
    

@cli.command()
@click.argument('doi')
@click.option('-r', '--remove', is_flag=True, help='remove')
def related(doi, remove=False):
    """Setup a watch for articles that are related to doi.
    """
           
    url = 'https://api.openalex.org/works'

    # We need an OpenAlex id
    f = f'doi:{doi}'

    data = get_data(url, params={'filter': f})
    if len(data['results']) == 0:
        richprint(f"Sorry, {doi} does not seem valid.")

    wid = data['results'][0]['id']

    if remove:
        c = db.execute('''delete from queries where filter = ?''',
                   (f'related_to:{wid}',))
        db.commit()
        richprint(f'Deleted {c.rowcount} rows')
    else:
        c = db.execute('''insert or ignore into queries(filter, description) values (?, ?)''',
               (f'related_to:{wid}', f'Related papers for {doi}'))
        
        db.commit()
        richprint(f'Added {c.rowcount} rows')        


#############
# Utilities #
#############
        
@cli.command()
@click.argument('sources', nargs=-1)
def bibtex(sources):
    """Generate bibtex entries for sources."""

    from .bibtex import dump_bibtex
    import json
    for source in sources:
        work, = db.execute('''select extra from sources where source = ?''', (source,)).fetchone()
        richprint(dump_bibtex(json.loads(work)))


@cli.command()
@click.argument('sources', nargs=-1)
def citation(sources):
    """Generate citation strings for sources."""

    from .bibtex import dump_bibtex
    import json
    for i, source in enumerate(sources):
        citation, = db.execute('''select json_extract(extra, '$.citation') from sources where source = ?''', (source,)).fetchone()
        richprint(f'{i + 1:2d}. {citation}')


@cli.command()
@click.argument('doi')
def unpaywall(doi):
    """
    """

    url =  f'https://api.unpaywall.org/v2/{doi}'
    params = {'email': config['openalex']['email']}

    data = requests.get(url, params).json()
    richprint(f'{data["title"]}, {data.get("journal_name") or ""}')
    richprint(f'Is open access: {data.get("is_oa", False)}')
    
    for loc in data.get('oa_locations', []):
        richprint(loc.get('url_for_pdf') or loc.get('url_for_landing_page'))


@cli.command()
def about():
    """Summary statistics of your db.
    """
    dbf = root / config['database']['db']
    richprint(f'Your database is located at {dbf}')
    kb = 1024
    mb = 1024 * kb
    gb = 1024 * mb
    richprint(f'Database size: {os.path.getsize(dbf) / gb:1.2f} GB')
    nsources, = db.execute('select count(source) from sources').fetchone()    
    richprint(f'You have {nsources} sources')

    
@cli.command()
@click.argument('sql')
def sql(sql):
    """Run the SQL command on the db.
    """
    for row in db.execute(sql).fetchall():
        richprint(row)


@cli.command()
@click.argument('source')
def visit(source):
    """Open source.
    """
    
    if source.startswith('http'):
        webbrowser.open(source, new=2)
    elif source.endswith('.pdf'):
        webbrowser.open(f'file://{source}')
    else:
        webbrowser.open(f'file://{source}')
    
        
if __name__ == '__main__':
    cli()
    
